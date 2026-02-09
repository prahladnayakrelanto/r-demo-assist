"""
PPT Slide Extractor
Extracts slides from PPTX files as PNG images for web display.
Also creates individual PPTX files for each slide for download.

Requirements:
    pip install python-pptx Pillow pdf2image

For Windows, also install poppler:
    1. Download from: https://github.com/osber/poppler-win/releases
    2. Extract and add bin folder to PATH

Alternative method using LibreOffice (if available):
    soffice --headless --convert-to pdf yourfile.pptx
    Then use pdf2image to convert PDF pages to images
"""

import os
import sys
import subprocess
import shutil
import copy
from pathlib import Path

# Configuration
INPUT_DIR = Path(".")
OUTPUT_DIR = Path("public/presentations/slides")
SLIDE_WIDTH = 1920  # Full HD width
SLIDE_HEIGHT = 1080  # Full HD height

def extract_with_libreoffice(pptx_path, output_folder):
    """Extract slides using LibreOffice (cross-platform)"""
    try:
        # First convert PPTX to PDF
        pdf_path = output_folder / f"{pptx_path.stem}.pdf"
        
        # Try different LibreOffice commands
        lo_commands = ['soffice', 'libreoffice', 'C:\\Program Files\\LibreOffice\\program\\soffice.exe']
        
        for lo_cmd in lo_commands:
            try:
                result = subprocess.run([
                    lo_cmd, '--headless', '--convert-to', 'pdf',
                    '--outdir', str(output_folder), str(pptx_path)
                ], capture_output=True, text=True, timeout=120)
                
                if result.returncode == 0:
                    print(f"  Converted to PDF using {lo_cmd}")
                    break
            except (FileNotFoundError, subprocess.TimeoutExpired):
                continue
        else:
            return False
        
        # Convert PDF to images using pdf2image
        try:
            from pdf2image import convert_from_path
            
            images = convert_from_path(
                str(output_folder / f"{pptx_path.stem}.pdf"),
                dpi=150,
                fmt='png'
            )
            
            for i, image in enumerate(images, 1):
                image_path = output_folder / f"slide_{i:03d}.png"
                image.save(str(image_path), 'PNG')
                print(f"  Saved slide {i}")
            
            # Clean up PDF
            (output_folder / f"{pptx_path.stem}.pdf").unlink(missing_ok=True)
            
            return True
            
        except ImportError:
            print("  pdf2image not installed. Run: pip install pdf2image")
            return False
            
    except Exception as e:
        print(f"  LibreOffice extraction failed: {e}")
        return False

def extract_with_comtypes(pptx_path, output_folder):
    """Extract slides using PowerPoint COM (Windows only)"""
    if sys.platform != 'win32':
        return False
    
    try:
        import comtypes.client
        
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1
        
        presentation = powerpoint.Presentations.Open(str(pptx_path.absolute()))
        
        for i, slide in enumerate(presentation.Slides, 1):
            image_path = output_folder / f"slide_{i:03d}.png"
            slide.Export(str(image_path.absolute()), "PNG", SLIDE_WIDTH, SLIDE_HEIGHT)
            print(f"  Saved slide {i}")
        
        presentation.Close()
        powerpoint.Quit()
        
        return True
        
    except ImportError:
        print("  comtypes not installed. Run: pip install comtypes")
        return False
    except Exception as e:
        print(f"  PowerPoint COM extraction failed: {e}")
        return False

def extract_text_from_slide(slide):
    """Extract all text content from a slide"""
    text_content = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip():
                        text_content.append(run.text.strip())
        
        # Handle tables
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text_content.append(cell.text.strip())
    
    return ' '.join(text_content)


def extract_slide_content(pptx_path, output_folder):
    """Extract text content from all slides for search functionality"""
    try:
        from pptx import Presentation
        
        src_prs = Presentation(str(pptx_path))
        slides_content = []
        
        print(f"  Extracting text content from {len(src_prs.slides)} slides...")
        
        for idx, slide in enumerate(src_prs.slides, 1):
            text = extract_text_from_slide(slide)
            slides_content.append({
                "slideNumber": idx,
                "content": text
            })
            print(f"    Slide {idx}: {len(text)} characters")
        
        # Save content to JSON
        import json
        content_path = output_folder / "content.json"
        with open(content_path, 'w', encoding='utf-8') as f:
            json.dump({"slides": slides_content}, f, indent=2, ensure_ascii=False)
        
        print(f"  Saved slide content to content.json")
        return slides_content
        
    except ImportError:
        print("  python-pptx not installed. Run: pip install python-pptx")
        return []
    except Exception as e:
        print(f"  Failed to extract slide content: {e}")
        return []


def extract_individual_slides_pptx(pptx_path, output_folder):
    """Extract each slide as a separate PPTX file using python-pptx"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        import copy
        
        # Open the source presentation
        src_prs = Presentation(str(pptx_path))
        slide_count = len(src_prs.slides)
        
        print(f"  Extracting {slide_count} slides as individual PPTX files...")
        
        for idx, slide in enumerate(src_prs.slides, 1):
            # Create a new presentation with the same slide size
            new_prs = Presentation()
            new_prs.slide_width = src_prs.slide_width
            new_prs.slide_height = src_prs.slide_height
            
            # Get the slide layout (use blank if available, otherwise first layout)
            try:
                blank_layout = new_prs.slide_layouts[6]  # Usually blank
            except:
                blank_layout = new_prs.slide_layouts[0]
            
            # Add a new slide
            new_slide = new_prs.slides.add_slide(blank_layout)
            
            # Copy all shapes from source slide to new slide
            for shape in slide.shapes:
                try:
                    # Clone the shape element
                    el = copy.deepcopy(shape.element)
                    new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
                except Exception as e:
                    # If cloning fails, try to recreate basic shapes
                    pass
            
            # Copy background if exists
            try:
                if slide.background.fill.type is not None:
                    new_slide.background.fill._fill = copy.deepcopy(slide.background.fill._fill)
            except:
                pass
            
            # Save the individual slide
            slide_pptx_path = output_folder / f"slide_{idx:03d}.pptx"
            new_prs.save(str(slide_pptx_path))
            print(f"    Created slide_{idx:03d}.pptx")
        
        return True
        
    except ImportError:
        print("  python-pptx not installed. Run: pip install python-pptx")
        return False
    except Exception as e:
        print(f"  Failed to extract individual PPTX slides: {e}")
        return False


def create_placeholder_slides(pptx_name, output_folder, num_slides=10):
    """Create placeholder slides when extraction fails"""
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("  Pillow not installed. Run: pip install Pillow")
        return False
    
    for i in range(1, num_slides + 1):
        # Create a slide image
        img = Image.new('RGB', (SLIDE_WIDTH, SLIDE_HEIGHT), color='#1F3459')
        draw = ImageDraw.Draw(img)
        
        # Try to use a nice font, fall back to default
        try:
            title_font = ImageFont.truetype("arial.ttf", 72)
            subtitle_font = ImageFont.truetype("arial.ttf", 36)
        except:
            title_font = ImageFont.load_default()
            subtitle_font = ImageFont.load_default()
        
        # Draw slide content
        title = f"Slide {i}"
        subtitle = pptx_name.replace('-', ' ').replace('.pptx', '')
        
        # Center the text
        draw.text((SLIDE_WIDTH//2, SLIDE_HEIGHT//2 - 50), title, 
                  fill='#FF3C00', font=title_font, anchor='mm')
        draw.text((SLIDE_WIDTH//2, SLIDE_HEIGHT//2 + 50), subtitle, 
                  fill='white', font=subtitle_font, anchor='mm')
        
        # Add slide number indicator
        draw.text((SLIDE_WIDTH - 100, SLIDE_HEIGHT - 50), f"{i}/{num_slides}", 
                  fill='#888888', font=subtitle_font, anchor='mm')
        
        image_path = output_folder / f"slide_{i:03d}.png"
        img.save(str(image_path), 'PNG')
        print(f"  Created placeholder slide {i}")
    
    return True

def process_pptx(pptx_path):
    """Process a single PPTX file"""
    print(f"\nProcessing: {pptx_path.name}")
    
    # Create output folder for this presentation
    folder_name = pptx_path.stem.replace(' ', '-').replace('&', 'and')
    output_folder = OUTPUT_DIR / folder_name
    output_folder.mkdir(parents=True, exist_ok=True)
    
    # Try different extraction methods
    success = False
    
    # Method 1: Try PowerPoint COM (Windows with PowerPoint installed)
    if not success:
        print("  Trying PowerPoint COM extraction...")
        success = extract_with_comtypes(pptx_path, output_folder)
    
    # Method 2: Try LibreOffice
    if not success:
        print("  Trying LibreOffice extraction...")
        success = extract_with_libreoffice(pptx_path, output_folder)
    
    # Method 3: Create placeholders
    if not success:
        print("  Creating placeholder slides...")
        success = create_placeholder_slides(pptx_path.name, output_folder)
    
    if success:
        # Extract text content from slides for search
        print("  Extracting slide content for search...")
        slides_content = extract_slide_content(pptx_path, output_folder)
        
        # Also extract individual slides as PPTX files for download
        print("  Extracting individual PPTX slides for download...")
        extract_individual_slides_pptx(pptx_path, output_folder)
        
        # Create metadata file
        slides = list(output_folder.glob("slide_*.png"))
        slide_pptx = list(output_folder.glob("slide_*.pptx"))
        metadata = {
            "name": pptx_path.stem,
            "folder": folder_name,
            "slideCount": len(slides),
            "slides": [s.name for s in sorted(slides)],
            "slidePptx": [s.name for s in sorted(slide_pptx)],
            "hasContent": len(slides_content) > 0
        }
        
        import json
        with open(output_folder / "metadata.json", 'w') as f:
            json.dump(metadata, f, indent=2)
        
        print(f"  Successfully processed {len(slides)} slides (images + PPTX + content)")
        return folder_name, len(slides)
    
    return None, 0

def main():
    """Main entry point"""
    print("=" * 60)
    print("PPT Slide Extractor")
    print("=" * 60)
    
    # Create output directory
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    
    # Find all PPTX files
    pptx_files = [
        INPUT_DIR / "Repository-Data & Analytics -Case-Studies.pptx",
        INPUT_DIR / "Repository-AI-Case-Studies 1.pptx"
    ]
    
    # Also look for any other PPTX files
    for pptx in INPUT_DIR.glob("*.pptx"):
        if pptx not in pptx_files:
            pptx_files.append(pptx)
    
    results = []
    for pptx_path in pptx_files:
        if pptx_path.exists():
            folder, count = process_pptx(pptx_path)
            if folder:
                results.append((folder, count))
    
    # Create index file for the React app
    if results:
        import json
        index = {
            "presentations": [
                {"folder": folder, "slideCount": count} 
                for folder, count in results
            ]
        }
        with open(OUTPUT_DIR / "index.json", 'w') as f:
            json.dump(index, f, indent=2)
        
        print("\n" + "=" * 60)
        print("Extraction complete!")
        print(f"Output directory: {OUTPUT_DIR.absolute()}")
        print("\nProcessed presentations:")
        for folder, count in results:
            print(f"  - {folder}: {count} slides")
        print("=" * 60)

if __name__ == "__main__":
    main()

